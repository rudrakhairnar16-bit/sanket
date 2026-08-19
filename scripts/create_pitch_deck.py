"""Yuva 6.0 Pitch Deck — 10 slides, V5 design, visual-first."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# V5 Tokens
BG = RGBColor(0x0B, 0x11, 0x20)
GOLD = RGBColor(0xC9, 0xA9, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xD0, 0xD0, 0xD0)
GRAY = RGBColor(0x88, 0x96, 0xAA)
DIM = RGBColor(0x50, 0x5E, 0x72)
CARD = RGBColor(0x12, 0x1C, 0x2E)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
BLUE = RGBColor(0x81, 0x8C, 0xF8)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
RED = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xC4, 0xB5, 0xFD)

blank = prs.slide_layouts[6]


def bg(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG


def txt(s, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.alignment = a
    return tb


def header(s, tag, title, sub=""):
    """V5-style spaced header."""
    txt(s, 0.8, 0.3, 12, 0.28, tag.upper(), sz=11, c=GOLD, b=True)
    # Gold line
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.62), Inches(1.5), Inches(0.018))
    ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
    txt(s, 0.8, 0.72, 12, 0.6, title, sz=32, c=WHITE, b=True)
    if sub:
        txt(s, 0.8, 1.3, 12, 0.35, sub, sz=13, c=GRAY)


def card(s, l, t, w, h, title="", lines=None, tc=GOLD, bc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD
    if bc:
        sh.line.color.rgb = bc; sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.1)
    if title:
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(12); p.font.color.rgb = tc; p.font.bold = True; p.space_after = Pt(5)
    if lines:
        start = 0 if title else 0
        for i, ln_text in enumerate(lines):
            p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
            p.text = ln_text
            p.font.size = Pt(11); p.font.color.rgb = LIGHT; p.space_after = Pt(2)
    return sh


def stat(s, l, t, w, h, num, desc, nc=GOLD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD
    sh.line.color.rgb = RGBColor(0x1E, 0x29, 0x3B); sh.line.width = Pt(1)
    tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.08)
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(32); p.font.color.rgb = nc; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(9); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER
    return sh


def footer(s):
    txt(s, 0.5, 7.05, 12.3, 0.3, "SANKET  ·  ISL FOR SARKARI CLERKS  ·  BEYOND WORDS  ·  KPGU", sz=8, c=DIM, a=PP_ALIGN.CENTER)


def arrow_row(s, items, y=2.3, colors=None):
    """Horizontal flow with arrows between items."""
    n = len(items)
    total_w = 11.7
    item_w = (total_w - 0.3 * (n - 1)) / n
    for i, item in enumerate(items):
        x = 0.8 + i * (item_w + 0.3)
        clr = colors[i] if colors else BLUE
        card(s, x, y, item_w, 0.9, "", [item], tc=clr)
        if i < n - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + item_w + 0.05), Inches(y + 0.3), Inches(0.2), Inches(0.25))
            arr.fill.solid(); arr.fill.fore_color.rgb = DIM; arr.line.fill.background()


# ==================== SLIDE 1: TITLE ====================
s = prs.slides.add_slide(blank); bg(s)

txt(s, 0, 0.6, 13.333, 0.3, "Y U V A  6 . 0  —  I N T E R - U N I V E R S I T Y  R O U N D", sz=11, c=GOLD, b=True, a=PP_ALIGN.CENTER)
txt(s, 0, 1.0, 13.333, 0.3, "PROBLEM STATEMENT :  I N D I A N   S I G N   L A N G U A G E   F O R   A L L", sz=12, c=GOLD, b=True, a=PP_ALIGN.CENTER)

ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.6), Inches(2.9), Inches(0.018))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()

txt(s, 0, 2.0, 13.333, 1.4, "SANKET", sz=88, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(s, 0, 3.5, 13.333, 0.5, "Indian Sign Language for Sarkari Clerks", sz=24, c=GOLD, a=PP_ALIGN.CENTER)
txt(s, 0, 4.1, 13.333, 0.4, '"From 30 days of training to 30 seconds of service."', sz=16, c=LIGHT, a=PP_ALIGN.CENTER)

ln2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(4.7), Inches(2.9), Inches(0.018))
ln2.fill.solid(); ln2.fill.fore_color.rgb = GOLD; ln2.line.fill.background()

txt(s, 0, 5.0, 13.333, 0.35, "TEAM BEYOND WORDS  ·  KPGU UNIVERSITY", sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(s, 0, 5.4, 13.333, 0.3, "Sector : Public Services   |   Audience : Government (B2G)   |   Track : Accessibility", sz=10, c=GOLD, a=PP_ALIGN.CENTER)

# 3 pill tags
tags = ["34 ISL Signs", "Live Demo Ready", "WebRTC Relay"]
for i, tg in enumerate(tags):
    x = 4.2 + i * 1.9
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.95), Inches(1.6), Inches(0.35))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD; sh.line.color.rgb = GOLD; sh.line.width = Pt(0.8)
    tf = sh.text_frame; p = tf.paragraphs[0]; p.text = tg; p.font.size = Pt(9); p.font.color.rgb = GOLD; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

footer(s)

# ==================== SLIDE 2: THE PROBLEM ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "The Problem", "The silence at the sarkari counter.", "18 million deaf citizens. 3.5 million clerks. Zero shared language.")

stat(s, 0.8, 2.0, 2.7, 1.5, "18M", "Deaf & hard-of-hearing\ncitizens in India", GOLD)
stat(s, 3.8, 2.0, 2.7, 1.5, "3.5M", "Government clerks in\nmunicipal offices", GOLD)
stat(s, 6.8, 2.0, 2.7, 1.5, "< 5%", "Have basic\nISL proficiency", RED)
stat(s, 9.8, 2.0, 2.7, 1.5, "0", "Tools built for\nthe counter itself", RED)

card(s, 0.8, 3.9, 5.7, 2.8, "ROOT CAUSE", [
    "ISL lessons exist off the counter —",
    "apps, courses, video libraries.",
    "",
    "Clerks never see them because",
    "the need is only visible AT the counter.",
    "",
    "No tool is installed where service happens.",
], tc=ORANGE, bc=ORANGE)

card(s, 6.9, 3.9, 5.7, 2.8, "CONSEQUENCE", [
    "Deaf citizens rely on written notes,",
    "family interpreters, or simply leave",
    "without service.",
    "",
    "No accountability. No recourse. No data.",
    "",
    "RPwD Act 2016 mandates accessibility.",
], tc=RED, bc=RED)

footer(s)

# ==================== SLIDE 3: OUR SOLUTION ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "The Solution", "One platform. Service, habit, governance.", "Three products. One counter. 30 seconds.")

# 3 pillar cards with icons
card(s, 0.8, 2.0, 3.7, 4.6, "1  ·  MOMENT", [
    "",
    "Sanket Sahayak",
    "Two-way service turn",
    "at the counter.",
    "",
    "→ App speaks first",
    "→ Citizen signs (camera)",
    "→ Clerk taps reply",
    "→ ISL chips + TTS",
    "",
    "30 seconds.",
    "That is the moment.",
], tc=GREEN)

card(s, 4.8, 2.0, 3.7, 4.6, "2  ·  HABIT", [
    "",
    "ISL Quest",
    "3-minute gamified lessons.",
    "XP. Streaks. Badges.",
    "",
    "→ Daily lesson on desk",
    "→ Webcam practice",
    "→ 34 municipal signs",
    "→ Spaced repetition",
    "",
    "The unit is 3 minutes,",
    "not a course.",
], tc=BLUE)

card(s, 8.8, 2.0, 3.7, 4.6, "3  ·  SCORE", [
    "",
    "Sugamya Score",
    "Departmental compliance.",
    "Real data. Not aspirational.",
    "",
    "45% Compliance",
    "30% Citizen Satisfaction",
    "15% Participation",
    "10% Human Safety Net",
    "",
    "Rewards honesty,",
    "not just completion.",
], tc=ORANGE)

footer(s)

# ==================== SLIDE 4: HOW IT WORKS ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "How It Works", "Boring, reliable technology.", "Next.js 14  +  MongoDB Atlas  +  MediaPipe  +  kNN  +  Vercel")

arrow_row(s, ["📷 MediaPipe\n21 landmarks", "🧠 kNN\nk=3, conf≥0.45", "💬 Clerk UI\nReact", "🔄 API\nServerless", "🗄️ MongoDB\nAtlas"], y=2.0)

card(s, 0.8, 3.3, 3.6, 1.8, "FRONTEND", [
    "Next.js 14 · Tailwind CSS",
    "PWA · Dark mode · EN / HI / MR",
    "Recharts · jsPDF · ServiceWorker",
])

card(s, 4.7, 3.3, 3.6, 1.8, "RECOGNITION", [
    "Euclidean kNN over 21 landmarks",
    "Wrist-relative normalization",
    "15-frame temporal smoothing",
    "On-device — no server round-trip",
], tc=ORANGE)

card(s, 8.6, 3.3, 3.9, 1.8, "ESCALATION", [
    "Low-confidence → Call interpreter",
    "WebSocket relay · text + ISL chips",
    "WebRTC video (Phase 2)",
    "Every relay logged → Score",
], tc=GREEN)

# Resilient data bar
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.0))
sh.fill.solid(); sh.fill.fore_color.rgb = CARD; sh.line.fill.background()
tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2)
p = tf.paragraphs[0]; p.text = "RESILIENT"; p.font.size = Pt(11); p.font.color.rgb = BLUE; p.font.bold = True
p2 = tf.add_paragraph(); p2.text = "Fail-fast DB (2s timeout) → mock fallback. Demo cannot crash. On-device ML → no internet dependency."; p2.font.size = Pt(11); p2.font.color.rgb = LIGHT

footer(s)

# ==================== SLIDE 5: POLICY ASK ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "Policy Recommendation", "Inclusive by law, not by charity.", "Mandate ISL training for every public-facing clerk.")

card(s, 0.8, 2.0, 5.7, 3.0, "THE RECOMMENDATION", [
    "",
    "Make ISL proficiency a mandatory",
    "annual requirement for all municipal",
    "public-facing staff.",
    "",
    "Delivered through a digital platform",
    "with real-time compliance tracking.",
    "",
    "Enforce via RPwD Act 2016 +",
    "Sugamya Bharat Abhiyan.",
], tc=GOLD, bc=GOLD)

card(s, 6.9, 2.0, 5.7, 3.0, "MEASURABLE KPIs", [
    "",
    "Daily Active Learners         > 70%",
    "30-day Streak                 > 50%",
    "Curriculum Completion         > 40%",
    "Citizen Satisfaction          > 4.0 / 5.0",
    "QR Scan Volume                1 / clerk / week",
    "ISL Champions                 10% of enrolled",
], tc=ORANGE)

# Cost bar
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(11.8), Inches(1.4))
sh.fill.solid(); sh.fill.fore_color.rgb = CARD; sh.line.color.rgb = GREEN; sh.line.width = Pt(1)
tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]; p.text = "COST BREAKDOWN  —  ₹2,50,000 / year per municipality (500 clerks)"; p.font.size = Pt(13); p.font.color.rgb = GREEN; p.font.bold = True; p.space_after = Pt(6)
items = [
    "Platform license (500 accounts)     ₹1,80,000     |     Curriculum + relay + dashboard     Included     |     Support     ₹70,000",
    "Per clerk per year     ₹500"
]
for item in items:
    p = tf.add_paragraph(); p.text = item; p.font.size = Pt(10); p.font.color.rgb = LIGHT; p.space_after = Pt(2)

footer(s)

# ==================== SLIDE 6: FEASIBILITY ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "Feasibility & Implementation", "Already built. Already working.", "Not a concept — a fully functional application deployed on Vercel.")

card(s, 0.8, 2.0, 5.7, 4.2, "WHAT'S READY NOW", [
    "",
    "✓  Next.js 14 app — 34 ISL signs",
    "✓  Live webcam tracking (MediaPipe)",
    "✓  kNN classifier — trainable per user",
    "✓  Human interpreter relay (simulated)",
    "✓  Gamification (XP, streaks, badges)",
    "✓  Admin dashboard + real analytics",
    "✓  Citizen QR feedback loop",
    "✓  Auth (JWT + bcrypt)",
    "✓  Mock DB fallback — zero crash",
    "✓  PWA offline mode",
], tc=GREEN)

card(s, 6.9, 2.0, 5.7, 1.6, "PILOT TIMELINE", [
    "Month 1 → Deploy in 1 municipality, 3 depts",
    "Month 2 → Onboard 50–100 clerks",
    "Month 3 → Measure DAL, iterate",
], tc=BLUE)

card(s, 6.9, 3.8, 5.7, 1.2, "FUTURE  ·  LIVE EMERGENCY INTERPRETER", [
    "Phase 2 → WebRTC video relay for low-confidence signs",
    "When AI can't read → trained human joins live via video call",
], tc=PURPLE, bc=PURPLE)

card(s, 6.9, 5.2, 5.7, 1.0, "TECHNICAL PROOF", [
    "Deploy Vercel · DB Atlas · ML on-device · Cost ₹500/clerk/yr",
], tc=ORANGE)

footer(s)

# ==================== SLIDE 7: SCALABILITY ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "Scalability Roadmap", "One office to the nation.", "Prove assists, not installs.")

phases = [
    ("PILOT", "Months 1–3", "1 municipality\n3 departments\n50–100 clerks", GREEN),
    ("CITY-WIDE", "Months 4–6", "All departments\n500+ clerks\n5–10 offices", BLUE),
    ("REGIONAL", "Months 7–12", "5 Tier 2/3 cities\n5,000+ clerks", ORANGE),
    ("STATE", "Year 2", "Gujarat + Maharashtra\n25,000+ certified", GOLD),
    ("NATIONAL", "Years 3–5", "36 states & UTs\n500,000+ trained", RED),
]
for i, (name, time, desc, clr) in enumerate(phases):
    x = 0.6 + i * 2.5
    card(s, x, 2.0, 2.2, 3.0, name, [time, "", desc], tc=clr)
    if i < 4:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.25), Inches(3.2), Inches(0.2), Inches(0.25))
        arr.fill.solid(); arr.fill.fore_color.rgb = DIM; arr.line.fill.background()

# Partners bar
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.2))
sh.fill.solid(); sh.fill.fore_color.rgb = CARD; sh.line.fill.background()
tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2)
p = tf.paragraphs[0]; p.text = "KEY PARTNERS"; p.font.size = Pt(11); p.font.color.rgb = GOLD; p.font.bold = True; p.space_after = Pt(4)
p2 = tf.add_paragraph(); p2.text = "GSRDM  ·  Maharashtra SRC  ·  DEPwD  ·  ISLRTC Certified  ·  ASHA-worker training model"; p2.font.size = Pt(10); p2.font.color.rgb = LIGHT

footer(s)

# ==================== SLIDE 8: IMPACT ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "Impact", "Everyone gains. Measurable at every level.", "Sanket creates tangible value for every stakeholder.")

card(s, 0.8, 2.0, 5.7, 1.5, "🧏  DEAF CITIZEN", [
    "Walk in. Be understood. No notes. No waiting.",
], tc=BLUE)

card(s, 6.9, 2.0, 5.7, 1.5, "🧑‍💼  CLERK", [
    "30-sec service. 3-min habit. Champion badge.",
], tc=GREEN)

card(s, 0.8, 3.8, 5.7, 1.5, "🏛️  MUNICIPALITY", [
    "Real-time compliance dashboard. Assisted citizen count.",
], tc=ORANGE)

card(s, 6.9, 3.8, 5.7, 1.5, "🇮🇳  NATION", [
    "ISL-trained workforce at every counter. Inclusive governance.",
], tc=GOLD)

# Impact number
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.6), Inches(9.3), Inches(1.0))
sh.fill.solid(); sh.fill.fore_color.rgb = CARD; sh.line.color.rgb = GOLD; sh.line.width = Pt(1)
tf = sh.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "18M deaf citizens  ×  3.5M government clerks  =  ONE PLATFORM"; p.font.size = Pt(18); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

footer(s)

# ==================== SLIDE 9: WHAT'S DIFFERENT ====================
s = prs.slides.add_slide(blank); bg(s)
header(s, "What's Different", "Built for the municipal counter, not the classroom.", "The first end-to-end ISL accessibility platform for Indian government clerks.")

card(s, 0.8, 2.0, 5.7, 2.3, "DIFFERENT BY DESIGN", [
    "",
    "→  Municipal vocabulary (Bill, Tax, Certificate)",
    "→  Clerk chat UI + ISL chips + TTS",
    "→  kNN — trainable per user",
    "→  Two-way: Sign → Text → ISL → Voice",
], tc=BLUE)

card(s, 6.9, 2.0, 5.7, 2.3, "WHAT MAKES US UNIQUE", [
    "",
    "→  Webcam + MediaPipe + kNN in-browser",
    "→  QR feedback loop for compliance",
    "→  Live human interpreter relay",
    "→  Offline-first PWA",
], tc=ORANGE)

# Quote
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(11.8), Inches(1.8))
sh.fill.solid(); sh.fill.fore_color.rgb = CARD; sh.line.color.rgb = GOLD; sh.line.width = Pt(1)
tf = sh.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.text = '"Not just an ASL translator. Not just a learning app.'
p.font.size = Pt(15); p.font.color.rgb = LIGHT; p.font.italic = True; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = 'The first end-to-end ISL accessibility platform for Indian government clerks —'
p2.font.size = Pt(15); p2.font.color.rgb = LIGHT; p2.font.italic = True; p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = '6 Lakh counters, where the counter IS the government."'
p3.font.size = Pt(15); p3.font.color.rgb = GOLD; p3.font.italic = True; p3.font.bold = True; p3.alignment = PP_ALIGN.CENTER

footer(s)

# ==================== SLIDE 10: THANK YOU ====================
s = prs.slides.add_slide(blank); bg(s)

txt(s, 0, 0.5, 13.333, 0.3, "T H A N K   Y O U", sz=12, c=GOLD, b=True, a=PP_ALIGN.CENTER)

ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.0), Inches(2.333), Inches(0.018))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()

txt(s, 0, 1.3, 13.333, 1.2, "SANKET", sz=72, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(s, 0, 2.7, 13.333, 0.4, "ISL for every sarkari counter.", sz=18, c=GRAY, a=PP_ALIGN.CENTER)

team = [("PJ", "Pratiksha Jawale", "Lead"), ("RK", "Rudra Khaire", "Dev"), ("MP", "Mahi Panchal", "Dev"), ("SP", "Suhani Pawar", "Dev"), ("SS", "Sheena Sharma", "Dev")]
for i, (ini, name, role) in enumerate(team):
    x = 1.8 + i * 2.0
    card(s, x, 3.5, 1.7, 1.8, ini, [name, role], tc=GOLD)

txt(s, 0, 5.6, 13.333, 0.35, "Dr. Kiran & Pallavi Patel Global University  ·  B.Tech CSE  ·  2nd Year", sz=11, c=GRAY, a=PP_ALIGN.CENTER)
txt(s, 0, 6.0, 13.333, 0.3, "Yuva 6.0 — Inter-University Round  |  Live Demo Ready", sz=10, c=GOLD, a=PP_ALIGN.CENTER)
txt(s, 0, 6.4, 13.333, 0.3, "sanket-isl.vercel.app", sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)

footer(s)

# Save
out = r"C:\Users\Rudra\Desktop\sanket\deliverables\Sanket-Pitch-Deck-Yuva6.pptx"
prs.save(out)
print(f"Saved: {out}")
