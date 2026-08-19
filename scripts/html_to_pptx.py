"""Convert presentation.html to PPTX using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
BLUE = RGBColor(0x81, 0x8C, 0xF8)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
PURPLE = RGBColor(0xC4, 0xB5, 0xFD)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
TAGLINE_COLOR = RGBColor(0x81, 0x8C, 0xF8)

blank_layout = prs.slide_layouts[6]  # blank


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=14, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def add_card(slide, left, top, width, height, title, body_lines, title_color=WHITE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(6)

    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.space_after = Pt(3)
    return shape


def add_stat_box(slide, left, top, width, height, number, desc):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(28)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    return shape


# ============ SLIDE 1: Title ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0, 0.8, 13.333, 0.5, "Yuva 6.0 — Inter-University Round", size=14, color=TAGLINE_COLOR, align=PP_ALIGN.CENTER)
add_text(slide, 0, 1.5, 13.333, 1.2, "Sanket", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0, 2.8, 13.333, 0.5, "ISL for Every Sarkari Counter", size=22, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 0, 3.8, 13.333, 0.5, "From 30 days of training, to 30 seconds of service.", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0, 4.5, 13.333, 0.5, "Not another recognizer — the only flow built for a government desk: sign in, reply in one tap, escalate to a human clerk live.", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 0, 5.3, 13.333, 0.5, "Team KPGU\nDr. Kiran & Pallavi Patel Global University · B.Tech CSE · 2nd Year", size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ============ SLIDE 2: Meet Vaishnavi ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Why this matters", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Meet Vaishnavi.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "She is deaf. Every water bill, complaint, and certificate in her city gets filed at a municipal counter — across from someone she cannot talk to.", size=14, color=GRAY)
add_card(slide, 0.8, 2.5, 3.5, 2.5, "Day 1", ["Vaishnavi walks in to pay her water bill.", "She signs 'bill'. The clerk does not understand.", "She writes a note. She waits."])
add_card(slide, 4.9, 2.5, 3.5, 2.5, "Day 5", ["The clerk learns 'wait', 'form' and 'name'", "in 3-minute lessons on the same desk.", "The counter starts to speak her language."])
add_card(slide, 9.0, 2.5, 3.5, 2.5, "Day 14", ["She signs 'thank you' and leaves with her", "certificate. No note. No wait.", "That is 30 seconds of service."], title_color=GREEN)

# ============ SLIDE 3: Problem ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "The Problem", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "A counter with no common language.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.8, "The moment in Vaishnavi's story is the same every single day, for millions of deaf citizens — because the people across the counter were never trained, and no tool lives where the counter is.", size=14, color=GRAY)

# Stats row
add_stat_box(slide, 0.8, 2.8, 2.8, 1.5, "18M", "Deaf & hard-of-hearing\ncitizens in India")
add_stat_box(slide, 3.9, 2.8, 2.8, 1.5, "3.5M", "Government clerks serving\nmunicipal offices")
add_stat_box(slide, 7.0, 2.8, 2.8, 1.5, "<5%", "Have basic\nISL proficiency")
add_stat_box(slide, 10.1, 2.8, 2.8, 1.5, "0", "Tools built for\nthe counter itself")

add_card(slide, 0.8, 4.8, 5.5, 1.8, "Root cause:", ["ISL lessons exist off the counter — apps, courses, video libraries.", "Clerks never see them because the need is only visible at the counter."], border_color=ORANGE)
add_card(slide, 6.9, 4.8, 5.5, 1.8, "Consequence:", ["Deaf citizens rely on written notes, family interpreters, or simply", "leave without service. No accountability, no recourse, no data."], border_color=INDIGO)

# ============ SLIDE 4: Sanket Sahayak ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "The Moment — Live in the demo", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Sanket Sahayak: serve a citizen in 30 seconds.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "One screen on the clerk's desk. The app speaks first, the citizen signs, the clerk taps a reply — text and voice go both ways.", size=14, color=GRAY)

# Flow boxes
flow_items = [
    ("App speaks first", '"Namaste, I am Sanket Sahayak"'),
    ("Citizen signs", "camera / one-tap demo"),
    ("Clerk sees text + hears it", "en · hi · mr"),
    ("Clerk taps one reply", "ISL chips + TTS to citizen"),
    ('"You helped a citizen"', "+25 XP · counter card"),
]
for i, (title, sub) in enumerate(flow_items):
    x = 0.8 + i * 2.5
    add_card(slide, x, 2.3, 2.2, 1.2, title, [sub], title_color=BLUE)

add_card(slide, 0.8, 4.0, 5.5, 1.2, "Why this wins adoption", ["The unit of engagement is 30 seconds, not a course.", "A clerk who won't attend training has no excuse to skip one tap."], title_color=GREEN)
add_card(slide, 6.9, 4.0, 5.5, 1.2, "Where the architecture shows", ["Fail-fast DB with mock fallback — counter never crashes.", "On-device recognition — no internet dependency."], title_color=INDIGO)

# Not a recognizer box
add_card(slide, 0.8, 5.5, 11.7, 1.5, "Not a recognizer — the desk is the differentiator", [
    "A recognizer: Sign → text, then what? Unknown sign → blank screen.",
    "Sanket: Low-confidence sign → one tap calls a live human interpreter. Every escalation is logged into the department score."
], title_color=PURPLE)

# ============ SLIDE 5: Solution ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Our Solution", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Moment → Habit → Score.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "The desk flow wins the clerk today; the gamified habit and the departmental score make it stick. One platform, three stages.", size=14, color=GRAY)

add_card(slide, 0.8, 2.5, 3.5, 3.0, "1 · Moment", ["Sanket Sahayak — a two-way service turn", "at the counter in under 30 seconds.", "This is what wins the clerk."], title_color=GREEN)
add_card(slide, 4.9, 2.5, 3.5, 3.0, "2 · Habit", ["ISL Quest — 3-minute gamified lessons.", "XP, streaks, badges, leaderboards", "and 34 municipal signs."], title_color=BLUE)
add_card(slide, 9.0, 2.5, 3.5, 3.0, "3 · Score", ["Sugamya Score — admin dashboard:", "assisted citizens, compliance,", "departmental leaderboards."], title_color=ORANGE)

# ============ SLIDE 6: Architecture ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Architecture", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Honest engineering.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "Built on modern, scalable infrastructure — Next.js 14 + MongoDB Atlas + MediaPipe + Vercel", size=14, color=GRAY)

# Data flow
flow_items2 = ["MediaPipe\n21 landmarks", "kNN Classifier\ndemo-grade", "Clerk UI\nReact", "API Routes\nserverless", "MongoDB Atlas\nUsers, modules"]
for i, item in enumerate(flow_items2):
    x = 0.8 + i * 2.5
    add_card(slide, x, 2.3, 2.2, 1.0, "", [item], title_color=BLUE)

add_card(slide, 0.8, 3.8, 3.5, 2.8, "Frontend", ["Next.js 14 App Router · Tailwind", "PWA · dark mode · EN/HI/MR", "Recharts · jsPDF · ServiceWorker"])
add_card(slide, 4.9, 3.8, 3.5, 2.8, "Recognition — honest", [
    "Euclidean kNN over MediaPipe landmarks",
    "wrist-relative normalization, 15-frame smoothing",
    "Runs on-device. The moat is what happens",
    "when it isn't enough."
], title_color=ORANGE)
add_card(slide, 9.0, 3.8, 3.5, 2.8, "Resilient data", [
    "Fail-fast connection (2s timeouts)",
    "If Mongo is down → mock data",
    "Demo cannot crash"
], title_color=INDIGO)

# ============ SLIDE 7: Policy ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Policy Recommendation", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Mandate ISL training for every public-facing clerk.", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "A simple, measurable, low-cost policy framework for municipal accessibility compliance.", size=14, color=GRAY)

add_card(slide, 0.8, 2.3, 5.5, 2.5, "Policy Proposal", [
    "Make ISL proficiency a mandatory annual requirement",
    "for all municipal public-facing staff, delivered",
    "through a digital platform with real-time tracking.",
    "",
    "70% DAL  |  50% streak  |  4.0/5 citizen score"
], title_color=INDIGO)

add_card(slide, 6.9, 2.3, 5.5, 2.5, "Measurable KPIs", [
    "Daily Active Learners (DAL)         >70%",
    "30-day Streak Achievement           >50%",
    "Curriculum Completion               >40%",
    "Citizen Satisfaction Score          >4.0 / 5.0",
    "QR Scan Volume                      1 / clerk / week",
    "ISL Champions                       10% of enrolled"
], title_color=ORANGE)

add_card(slide, 0.8, 5.2, 11.7, 1.8, "Cost Breakdown — ₹2,50,000 / year per municipality (up to 500 clerks)", [
    "Platform license (500 clerk accounts)              ₹1,80,000",
    "Self-paced ISL curriculum (25+ municipal signs)    Included",
    "Human interpreter relay — deaf ↔ trained clerk     Included",
    "Admin dashboard, analytics & QR feedback           Included",
    "Technical support & maintenance                    ₹70,000",
    "Total per municipality                             ₹2,50,000"
], title_color=GREEN)

# ============ SLIDE 8: Feasibility ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Feasibility & Implementation", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Already built. Already working.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "Not a concept — a fully functional application. Deployed on Vercel with MongoDB Atlas.", size=14, color=GRAY)

add_card(slide, 0.8, 2.3, 5.5, 3.5, "What's Ready", [
    "✓  Next.js 14 app with 25 ISL signs",
    "✓  Live webcam hand tracking (MediaPipe)",
    "✓  Human interpreter relay — deaf ↔ trained clerk",
    "✓  Gamification (XP, streaks, badges, levels)",
    "✓  Admin dashboard with real analytics",
    "✓  Citizen QR feedback loop",
    "✓  Auth (JWT + bcrypt + httpOnly cookies)",
    "✓  Mock DB fallback — zero crash risk"
], title_color=GREEN)

add_card(slide, 6.9, 2.3, 5.5, 3.5, "Pilot Timeline", [
    "Month 1: Deploy pilot in 1 municipality, 3 departments",
    "Month 2: Onboard & train 50-100 clerks",
    "Month 3: Measure DAL, collect QR feedback, iterate",
    "",
    "Phase 2: Live human interpreter relay via WebRTC",
    "— emergency escalation when AI can't read a sign",
    "",
    "Deploy: Vercel · serverless",
    "DB: MongoDB Atlas · real data",
    "Cost: ₹2.5L/yr · ₹500/clerk/yr"
], title_color=INDIGO)

# ============ SLIDE 9: Roadmap ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Scalability Roadmap", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "From one municipal office to 500,000 trained staff.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "A phased rollout designed for sustainable, measurable growth.", size=14, color=GRAY)

phases = [
    ("Pilot", "1 municipality · 3 depts", "50-100 clerks", "Months 1-3"),
    ("City-Wide", "All departments", "500+ clerks", "Months 4-6"),
    ("Regional", "5 Tier 2/3 cities", "5,000+ learners", "Months 7-12"),
    ("State", "Gujarat · Maharashtra", "25,000+ certified", "Year 2"),
    ("National", "36 states & UTs", "500,000+ trained", "Years 3-5"),
]
for i, (name, scope, target, time) in enumerate(phases):
    x = 0.8 + i * 2.5
    add_card(slide, x, 2.3, 2.2, 2.0, name, [scope, target, time], title_color=BLUE)

# ============ SLIDE 10: ISL Gap ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "We know the gaps — #1", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "ISL is not one language.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "Signs vary by region, community, and age. A single national recognizer is the wrong target — so we don't build one.", size=14, color=GRAY)

add_card(slide, 0.8, 2.5, 5.5, 3.0, "The honest problem", [
    "→  Same word, different handshape in Gujarat vs Maharashtra",
    "→  School ISL ≠ street ISL ≠ home signs",
    "→  A model trained on one community fails the next town"
], title_color=ORANGE)

add_card(slide, 6.9, 2.5, 5.5, 3.0, "How we design for it", [
    "→  ISLRTC-sourced vocabulary + icons + text + TTS",
    "→  Regional variant packs, deployable per municipality",
    "→  60-second on-device calibration per clerk",
    "→  Low-confidence sign → human interpreter relay"
], title_color=GREEN)

# ============ SLIDE 11: Rural ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "We know the gaps — #2", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Rural counters are the hardest — and the most important.", size=32, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "The taluka office is where the counter IS the government. It also has the worst power, internet, and training.", size=14, color=GRAY)

add_card(slide, 0.8, 2.5, 3.5, 2.5, "Power & internet", ["Offline-first PWA: core flow runs", "with no network after first load.", "MediaPipe runs on-device."])
add_card(slide, 4.9, 2.5, 3.5, 2.5, "No trainer on site", ["Train-the-trainer: one ISL-trained", "talati per circle becomes the", "local node — ASHA-worker model."])
add_card(slide, 9.0, 2.5, 3.5, 2.5, "Adoption, not installs", ["Pilot = one taluka office,", "measured on assisted citizens.", "If it works worst, it works everywhere."])

# ============ SLIDE 12: Impact ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Impact", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Everyone gains. Measurable at every level.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "Sanket creates tangible value for every stakeholder in the ecosystem.", size=14, color=GRAY)

add_card(slide, 0.8, 2.3, 5.5, 1.8, "Deaf Citizen", ["Walk into any govt office and be understood.", "No more written notes or turned away."], title_color=BLUE)
add_card(slide, 6.9, 2.3, 5.5, 1.8, "Clerk", ["30-second service today, 3-min/day habit.", "Champion badge for career recognition."], title_color=GREEN)
add_card(slide, 0.8, 4.5, 5.5, 1.8, "Municipality", ["Measurable accessibility compliance.", "Real-time dashboard. Assisted-citizen count."], title_color=ORANGE)
add_card(slide, 6.9, 4.5, 5.5, 1.8, "Nation", ["ISL-trained workforce reaches every counter.", "Truly inclusive governance at scale."], title_color=PURPLE)

add_card(slide, 0.8, 6.5, 11.7, 0.8, "", ["18 million deaf citizens × 3.5 million government clerks = one platform to bridge the gap."], title_color=BLUE)

# ============ SLIDE 13: Differentiation ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0.8, 0.5, 5, 0.4, "Differentiation", size=12, color=TAGLINE_COLOR)
add_text(slide, 0.8, 0.9, 11, 0.6, "Built for the municipal counter, not the classroom.", size=36, color=WHITE, bold=True)
add_text(slide, 0.8, 1.5, 11, 0.5, "What sets Sanket apart from every other accessibility product.", size=14, color=GRAY)

add_card(slide, 0.8, 2.3, 5.5, 3.0, "Different by design", [
    "→  Municipal-specific vocabulary (Bill, Tax, Certificate)",
    "→  Clerk-facing chat UI with ISL symbol chips + TTS",
    "→  kNN classifier — trainable per user",
    "→  Two-way: Sign → Text → ISL Symbols → TTS"
], title_color=INDIGO)

add_card(slide, 6.9, 2.3, 5.5, 3.0, "What makes us unique", [
    "→  Real-time webcam + MediaPipe + kNN in-browser",
    "→  Citizen feedback QR loop for compliance",
    "→  WhatsApp nudge for low-engagement clerks",
    "→  Offline-first PWA for low-connectivity"
], title_color=ORANGE)

add_card(slide, 0.8, 5.8, 11.7, 1.0, "", [
    '"Not just an ASL translator. Not just a learning app. The first end-to-end ISL accessibility platform for Indian government clerks."'
], title_color=PURPLE)

# ============ SLIDE 14: Thank You ============
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_text(slide, 0, 1.0, 13.333, 0.4, "Thank You", size=14, color=TAGLINE_COLOR, align=PP_ALIGN.CENTER)
add_text(slide, 0, 1.5, 13.333, 1.0, "Sanket", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0, 2.6, 13.333, 0.5, "ISL for every sarkari counter.", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Team
team = [
    ("PJ", "Pratiksha Jawale", "Team Lead"),
    ("RK", "Rudra Khaire", "Developer"),
    ("MP", "Mahi Panchal", "Developer"),
    ("SP", "Suhani Pawar", "Developer"),
    ("SS", "Sheena Sharma", "Developer"),
]
for i, (initials, name, role) in enumerate(team):
    x = 2.5 + i * 1.8
    add_card(slide, x, 3.5, 1.5, 1.5, initials, [name, role], title_color=BLUE)

add_text(slide, 0, 5.3, 13.333, 0.5, "Dr. Kiran & Pallavi Patel Global University · B.Tech CSE · 2nd Year", size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 0, 5.8, 13.333, 0.4, "Yuva 6.0 — Inter-University Round  |  Live Demo", size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 0, 6.3, 13.333, 0.4, "sanket-isl.vercel.app", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Save
output_path = r"C:\Users\Rudra\Desktop\sanket\Sanket-Presentation-Final.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
